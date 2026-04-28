"""
ICC Companion PowerPoint Generator
Creates a professional 12-slide presentation matching the website design
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# ICC Companion Color Palette (from CSS)
PRIMARY = RGBColor(37, 99, 235)      # #2563eb
PRIMARY_DARK = RGBColor(29, 78, 216)  # #1d4ed8
SECONDARY = RGBColor(16, 185, 129)    # #10b981
WARNING = RGBColor(245, 158, 11)      # #f59e0b
LIGHT = RGBColor(248, 250, 252)       # #f8fafc
DARK = RGBColor(30, 41, 59)           # #1e293b
GRAY = RGBColor(100, 116, 139)        # #64748b
WHITE = RGBColor(255, 255, 255)

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_background_gradient(slide1, PRIMARY, PRIMARY_DARK)
    
    # Logo
    title_box = slide1.shapes.add_textbox(Inches(0), Inches(1.5), Inches(10), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "🎓"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(72)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Main Title
    title_box2 = slide1.shapes.add_textbox(Inches(0), Inches(2.5), Inches(10), Inches(1))
    title_frame2 = title_box2.text_frame
    title_frame2.text = "ICC COMPANION"
    title_para2 = title_frame2.paragraphs[0]
    title_para2.font.size = Pt(60)
    title_para2.font.bold = True
    title_para2.font.color.rgb = WHITE
    title_para2.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide1.shapes.add_textbox(Inches(0), Inches(3.5), Inches(10), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Icon Commerce College Campus Portal"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = WHITE
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Tagline
    tagline_box = slide1.shapes.add_textbox(Inches(0), Inches(4.5), Inches(10), Inches(0.6))
    tagline_frame = tagline_box.text_frame
    tagline_frame.text = "Your All-in-One College Management Solution"
    tagline_para = tagline_frame.paragraphs[0]
    tagline_para.font.size = Pt(24)
    tagline_para.font.color.rgb = WHITE
    tagline_para.alignment = PP_ALIGN.CENTER
    
    # Footer
    footer_box = slide1.shapes.add_textbox(Inches(0), Inches(6), Inches(10), Inches(0.5))
    footer_frame = footer_box.text_frame
    footer_frame.text = "Final Year Project Presentation"
    footer_para = footer_frame.paragraphs[0]
    footer_para.font.size = Pt(18)
    footer_para.font.color.rgb = WHITE
    footer_para.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Project Overview
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_background_solid(slide2, LIGHT)
    
    add_title(slide2, "📋 Project Overview", DARK)
    
    content = slide2.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(1))
    content_frame = content.text_frame
    content_frame.text = "A professional web-based management system designed to digitize and simplify campus life for both students and staff."
    content_para = content_frame.paragraphs[0]
    content_para.font.size = Pt(22)
    content_para.font.color.rgb = DARK
    content_para.alignment = PP_ALIGN.CENTER
    
    # Three cards
    cards_data = [
        ("🎯", "Mission", "Streamline academic operations and enhance communication"),
        ("⚡", "Real-Time", "Instant access to attendance, schedules, and announcements"),
        ("📱", "Accessible", "Responsive design works on all devices")
    ]
    
    x_start = 0.5
    for i, (icon, title, desc) in enumerate(cards_data):
        x = x_start + (i * 3.2)
        add_card(slide2, x, 3.2, 2.8, 2.5, icon, title, desc)
    
    # Slide 3: Tech Stack
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_background_solid(slide3, WHITE)
    
    add_title(slide3, "🛠️ Technology Stack", DARK)
    
    subtitle3 = slide3.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(0.5))
    subtitle3_frame = subtitle3.text_frame
    subtitle3_frame.text = "Built with proven and reliable technologies"
    subtitle3_para = subtitle3_frame.paragraphs[0]
    subtitle3_para.font.size = Pt(20)
    subtitle3_para.font.color.rgb = GRAY
    
    # Frontend
    frontend_box = slide3.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(4.5), Inches(2.5))
    frontend_frame = frontend_box.text_frame
    frontend_frame.text = "Frontend\n\n• HTML5 & CSS3 - Semantic structure and modern responsive design\n• Vanilla JavaScript - Dynamic DOM manipulation and async API handling"
    format_text_box(frontend_frame, 20, DARK)
    
    # Backend
    backend_box = slide3.shapes.add_textbox(Inches(5.2), Inches(2.5), Inches(4.5), Inches(2.5))
    backend_frame = backend_box.text_frame
    backend_frame.text = "Backend\n\n• PHP - Server-side logic and API endpoints\n• MySQL - Relational database management\n• Apache - Web server via XAMPP"
    format_text_box(backend_frame, 20, DARK)
    
    # Tech badges
    badges = ["HTML5", "CSS3", "JavaScript", "PHP", "MySQL", "JSON API"]
    badge_y = 5.5
    badge_x_start = 1.5
    for i, badge in enumerate(badges):
        x = badge_x_start + (i % 3) * 2.5
        y = badge_y + (i // 3) * 0.6
        add_badge(slide3, x, y, badge)
    
    # Slide 4: Student Features
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    set_background_solid(slide4, LIGHT)
    
    add_title(slide4, "👨‍🎓 Student Dashboard Features", DARK)
    
    features = [
        "Attendance Tracker - Visual progress bars with automatic eligibility calculations",
        "Exam Timetable - Real-time schedules with countdown timers",
        "Announcements - Centralized feed for college notices and events",
        "Lost & Found Portal - Browse and report lost/found items",
        "Academic Resources - Access to class routines and exam schedules"
    ]
    
    add_bullet_list(slide4, 2, features, 20)
    
    # Slide 5: Admin Features
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    set_background_solid(slide5, WHITE)
    
    add_title(slide5, "👩‍💼 Admin Management Panel", DARK)
    
    admin_features = [
        "Student Information System - Full CRUD operations for student profiles",
        "Smart Attendance Marker - Quick checkbox-based interface for marking attendance",
        "Exam Controller - Manage academic calendar and schedule exams",
        "Global Communications - Post and manage announcements",
        "Lost & Found Manager - Review and manage reported items",
        "Dynamic File Manager - Upload and update routines and schedules"
    ]
    
    add_bullet_list(slide5, 2, admin_features, 20)
    
    # Slide 6: Key Statistics
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    set_background_gradient(slide6, SECONDARY, RGBColor(5, 150, 105))
    
    add_title(slide6, "📊 System Capabilities", WHITE)
    
    stats = [
        ("8", "Core Modules"),
        ("52", "API Endpoints"),
        ("9", "Database Tables"),
        ("2", "User Roles"),
        ("100%", "Mobile Responsive")
    ]
    
    # First row - 3 stats
    for i in range(3):
        x = 1 + (i * 2.8)
        add_stat_box(slide6, x, 2.5, stats[i][0], stats[i][1])
    
    # Second row - 2 stats
    for i in range(2):
        x = 2.4 + (i * 2.8)
        add_stat_box(slide6, x, 4.5, stats[i+3][0], stats[i+3][1])
    
    # Slide 7: Database Schema
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    set_background_solid(slide7, LIGHT)
    
    add_title(slide7, "💾 Database Architecture", DARK)
    
    db_subtitle = slide7.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(0.5))
    db_subtitle_frame = db_subtitle.text_frame
    db_subtitle_frame.text = "Highly structured relational database: campus_portal"
    db_subtitle_para = db_subtitle_frame.paragraphs[0]
    db_subtitle_para.font.size = Pt(20)
    db_subtitle_para.font.color.rgb = DARK
    
    # Table data
    table_data = [
        ["Table", "Purpose", "Key Relations"],
        ["departments", "Academic divisions", "Parent to subjects, students"],
        ["students", "Student profiles & credentials", "Links to attendance, dept"],
        ["attendance", "Daily class records", "Student + Subject FK"],
        ["exams", "Exam schedules", "Subject FK"],
        ["announcements", "System-wide notices", "Department targeting"]
    ]
    
    add_table(slide7, 0.5, 2.5, 9, 3.5, table_data)
    
    # Slide 8: System Architecture
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    set_background_solid(slide8, WHITE)
    
    add_title(slide8, "🏗️ System Architecture", DARK)
    
    arch_items = [
        "👤 User (Student/Admin)",
        "↓",
        "🖥️ Web Interface (HTML/CSS/JS)",
        "↓",
        "🔌 API Layer (PHP Endpoints)",
        "↓",
        "💾 MySQL Database + 📁 File Storage"
    ]
    
    y_pos = 2.5
    for item in arch_items:
        arch_box = slide8.shapes.add_textbox(Inches(2), Inches(y_pos), Inches(6), Inches(0.5))
        arch_frame = arch_box.text_frame
        arch_frame.text = item
        arch_para = arch_frame.paragraphs[0]
        arch_para.font.size = Pt(24 if item != "↓" else 32)
        arch_para.font.bold = True if item != "↓" else False
        arch_para.font.color.rgb = PRIMARY if item == "↓" else DARK
        arch_para.alignment = PP_ALIGN.CENTER
        y_pos += 0.6
    
    # Slide 9: Technical Highlights
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    set_background_solid(slide9, LIGHT)
    
    add_title(slide9, "🔬 Technical Implementation", DARK)
    
    tech_cards = [
        ("🔒 Security", "Prepared Statements - SQL injection prevention\nSession Management - Secure PHP sessions"),
        ("⚡ Performance", "AJAX/Fetch API - Asynchronous data loading\nJSON Exchange - Lightweight data transfer"),
    ]
    
    for i, (title, desc) in enumerate(tech_cards):
        x = 0.5 + (i * 4.8)
        add_tech_card(slide9, x, 2.5, 4.5, 1.8, title, desc)
    
    # Data Integrity card (full width)
    add_tech_card(slide9, 0.5, 4.5, 9, 1.5, "🎯 Data Integrity", 
                  "Database Transactions - ACID properties ensure atomic operations")
    
    # Slide 10: Development Process
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    set_background_solid(slide10, WHITE)
    
    add_title(slide10, "📅 Development Roadmap", DARK)
    
    dev_subtitle = slide10.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(0.5))
    dev_subtitle_frame = dev_subtitle.text_frame
    dev_subtitle_frame.text = "Structured 30-day development cycle"
    dev_subtitle_para = dev_subtitle_frame.paragraphs[0]
    dev_subtitle_para.font.size = Pt(20)
    dev_subtitle_para.font.color.rgb = GRAY
    
    weeks = [
        ("Week 1: Foundation", "Core architecture, database setup, and authentication system"),
        ("Week 2: Student Experience", "Attendance tracker, timetable, and Lost & Found modules"),
        ("Week 3: Admin Control", "Management tools for students, attendance, and exams"),
        ("Week 4: Polish & Deploy", "Announcement system, UI/UX refinements, and testing")
    ]
    
    for i, (title, desc) in enumerate(weeks):
        x = 0.5 + (i % 2) * 4.8
        y = 2.8 + (i // 2) * 2.2
        add_week_card(slide10, x, y, 4.5, 1.8, title, desc)
    
    # Slide 11: Key Achievements
    slide11 = prs.slides.add_slide(prs.slide_layouts[6])
    set_background_gradient(slide11, WARNING, RGBColor(217, 119, 6))
    
    add_title(slide11, "🏆 Key Achievements", WHITE)
    
    achievements = [
        "Fully functional dual-role system (Student & Admin)",
        "Real-time attendance tracking with automatic calculations",
        "Secure authentication and SQL injection prevention",
        "Mobile-responsive design for all devices",
        "Complete CRUD operations across all modules",
        "Professional-grade code architecture"
    ]
    
    add_bullet_list_white(slide11, 2.5, achievements, 22)
    
    # Slide 12: Thank You
    slide12 = prs.slides.add_slide(prs.slide_layouts[6])
    set_background_gradient(slide12, PRIMARY, PRIMARY_DARK)
    
    # Logo
    logo_box = slide12.shapes.add_textbox(Inches(0), Inches(1.5), Inches(10), Inches(1))
    logo_frame = logo_box.text_frame
    logo_frame.text = "🎓"
    logo_para = logo_frame.paragraphs[0]
    logo_para.font.size = Pt(72)
    logo_para.alignment = PP_ALIGN.CENTER
    
    # Thank You
    thanks_box = slide12.shapes.add_textbox(Inches(0), Inches(2.5), Inches(10), Inches(1))
    thanks_frame = thanks_box.text_frame
    thanks_frame.text = "Thank You!"
    thanks_para = thanks_frame.paragraphs[0]
    thanks_para.font.size = Pt(60)
    thanks_para.font.bold = True
    thanks_para.font.color.rgb = WHITE
    thanks_para.alignment = PP_ALIGN.CENTER
    
    # Project name
    project_box = slide12.shapes.add_textbox(Inches(0), Inches(3.8), Inches(10), Inches(0.6))
    project_frame = project_box.text_frame
    project_frame.text = "ICC COMPANION"
    project_para = project_frame.paragraphs[0]
    project_para.font.size = Pt(32)
    project_para.font.bold = True
    project_para.font.color.rgb = WHITE
    project_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide12.shapes.add_textbox(Inches(0), Inches(4.5), Inches(10), Inches(0.5))
    sub_frame = sub_box.text_frame
    sub_frame.text = "Icon Commerce College Campus Portal"
    sub_para = sub_frame.paragraphs[0]
    sub_para.font.size = Pt(24)
    sub_para.font.color.rgb = WHITE
    sub_para.alignment = PP_ALIGN.CENTER
    
    # Q&A
    qa_box = slide12.shapes.add_textbox(Inches(0), Inches(5.5), Inches(10), Inches(0.5))
    qa_frame = qa_box.text_frame
    qa_frame.text = "Questions & Answers"
    qa_para = qa_frame.paragraphs[0]
    qa_para.font.size = Pt(28)
    qa_para.font.color.rgb = WHITE
    qa_para.alignment = PP_ALIGN.CENTER
    
    return prs

def set_background_solid(slide, color):
    """Set solid background color"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def set_background_gradient(slide, color1, color2):
    """Set gradient background"""
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 135
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[1].color.rgb = color2

def add_title(slide, text, color):
    """Add title to slide"""
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = text
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = color

def add_card(slide, x, y, width, height, icon, title, description):
    """Add a feature card"""
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(226, 232, 240)
    
    text_frame = shape.text_frame
    text_frame.clear()
    
    # Icon
    p1 = text_frame.paragraphs[0]
    p1.text = icon
    p1.font.size = Pt(40)
    p1.alignment = PP_ALIGN.CENTER
    
    # Title
    p2 = text_frame.add_paragraph()
    p2.text = title
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = PRIMARY
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(10)
    
    # Description
    p3 = text_frame.add_paragraph()
    p3.text = description
    p3.font.size = Pt(14)
    p3.font.color.rgb = GRAY
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(5)

def add_badge(slide, x, y, text):
    """Add a technology badge"""
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(2), Inches(0.4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    shape.line.color.rgb = PRIMARY
    
    text_frame = shape.text_frame
    text_frame.text = text
    p = text_frame.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

def add_bullet_list(slide, y_start, items, font_size):
    """Add bullet list"""
    list_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_start), Inches(8.5), Inches(4.5))
    text_frame = list_box.text_frame
    
    for i, item in enumerate(items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = "✓ " + item
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK
        p.space_after = Pt(15)
        p.level = 0

def add_bullet_list_white(slide, y_start, items, font_size):
    """Add bullet list with white text"""
    list_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_start), Inches(8.5), Inches(4))
    text_frame = list_box.text_frame
    
    for i, item in enumerate(items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = "✓ " + item
        p.font.size = Pt(font_size)
        p.font.color.rgb = WHITE
        p.space_after = Pt(15)
        p.level = 0

def add_stat_box(slide, x, y, number, label):
    """Add statistics box"""
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(2.3), Inches(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255, alpha=25)
    shape.line.color.rgb = WHITE
    
    text_frame = shape.text_frame
    text_frame.clear()
    
    # Number
    p1 = text_frame.paragraphs[0]
    p1.text = number
    p1.font.size = Pt(44)
    p1.font.bold = True
    p1.font.color.rgb = WHITE
    p1.alignment = PP_ALIGN.CENTER
    
    # Label
    p2 = text_frame.add_paragraph()
    p2.text = label
    p2.font.size = Pt(18)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(5)

def add_table(slide, x, y, width, height, data):
    """Add table to slide"""
    rows = len(data)
    cols = len(data[0])
    
    table = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(width), Inches(height)).table
    
    # Header row
    for col_idx, cell_text in enumerate(data[0]):
        cell = table.cell(0, col_idx)
        cell.text = cell_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY
        
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(16)
        paragraph.font.bold = True
        paragraph.font.color.rgb = WHITE
    
    # Data rows
    for row_idx in range(1, rows):
        for col_idx, cell_text in enumerate(data[row_idx]):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text
            
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(14)
            paragraph.font.color.rgb = DARK

def add_tech_card(slide, x, y, width, height, title, description):
    """Add technical feature card"""
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(226, 232, 240)
    
    text_frame = shape.text_frame
    text_frame.clear()
    
    # Title
    p1 = text_frame.paragraphs[0]
    p1.text = title
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = PRIMARY
    
    # Description
    p2 = text_frame.add_paragraph()
    p2.text = description
    p2.font.size = Pt(16)
    p2.font.color.rgb = GRAY
    p2.space_before = Pt(10)

def add_week_card(slide, x, y, width, height, title, description):
    """Add development week card"""
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(226, 232, 240)
    
    text_frame = shape.text_frame
    text_frame.clear()
    
    # Title
    p1 = text_frame.paragraphs[0]
    p1.text = title
    p1.font.size = Pt(18)
    p1.font.bold = True
    p1.font.color.rgb = PRIMARY
    
    # Description
    p2 = text_frame.add_paragraph()
    p2.text = description
    p2.font.size = Pt(14)
    p2.font.color.rgb = GRAY
    p2.space_before = Pt(8)

def format_text_box(text_frame, font_size, color):
    """Format text box"""
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = color
        paragraph.space_after = Pt(10)

if __name__ == "__main__":
    print("Creating ICC Companion PowerPoint Presentation...")
    prs = create_presentation()
    prs.save("ICC_Companion_Presentation.pptx")
    print("✓ Presentation created successfully: ICC_Companion_Presentation.pptx")
    print("\nPresentation includes 12 slides:")
    print("1. Title Slide")
    print("2. Project Overview")
    print("3. Technology Stack")
    print("4. Student Dashboard Features")
    print("5. Admin Management Panel")
    print("6. System Capabilities")
    print("7. Database Architecture")
    print("8. System Architecture")
    print("9. Technical Implementation")
    print("10. Development Roadmap")
    print("11. Key Achievements")
    print("12. Thank You")
