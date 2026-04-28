from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
PRIMARY = RGBColor(0x25, 0x63, 0xEB)
SECONDARY = RGBColor(0x10, 0xB9, 0x81)
WARNING = RGBColor(0xF5, 0x9E, 0x0B)
DANGER = RGBColor(0xEF, 0x44, 0x44)
DARK = RGBColor(0x1E, 0x29, 0x3B)
GRAY = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF8, 0xFA, 0xFC)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name='Segoe UI'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_card(slide, left, top, width, height, title, body_lines, icon="", title_color=DARK, accent_color=None):
    card = add_rect(slide, left, top, width, height, WHITE, BORDER)
    y = top + 0.25
    if icon:
        add_text_box(slide, left + 0.25, y, 0.5, 0.4, icon, font_size=20)
        y += 0.4
    add_text_box(slide, left + 0.25, y, width - 0.5, 0.35, title, font_size=16, bold=True, color=title_color)
    y += 0.4
    for line in body_lines:
        add_text_box(slide, left + 0.25, y, width - 0.5, 0.28, line, font_size=11, color=GRAY)
        y += 0.26
    if accent_color:
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Emu(Inches(0.04)))
        accent.fill.solid()
        accent.fill.fore_color.rgb = accent_color
        accent.line.fill.background()

def slide_header(slide, title, dark=False):
    color = WHITE if dark else PRIMARY
    add_text_box(slide, 0.8, 0.5, 10, 0.6, title, font_size=28, bold=True, color=color)
    line_color = RGBColor(0x33,0x41,0x55) if dark else BORDER
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(11.7), Emu(Inches(0.02)))
    ln.fill.solid()
    ln.fill.fore_color.rgb = line_color
    ln.line.fill.background()

# ── SLIDE 1: Title ──
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, PRIMARY)
add_text_box(s, 1, 1.0, 11.3, 0.5, "Final Year Project Presentation", font_size=20, color=RGBColor(0xBA,0xE6,0xFD), alignment=PP_ALIGN.CENTER)
add_text_box(s, 1, 1.6, 11.3, 1.0, "ICC Companion", font_size=52, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(s, 1, 2.6, 11.3, 0.5, "A Web-Based Student Portal & College Management System", font_size=18, color=RGBColor(0xBA,0xE6,0xFD), alignment=PP_ALIGN.CENTER)

info = add_rect(s, 3.5, 3.5, 6.3, 1.8, RGBColor(0x1D,0x4E,0xD8), RGBColor(0x60,0x9F,0xFF))
add_text_box(s, 3.8, 3.6, 2.5, 0.3, "SUBMITTED BY", font_size=10, color=RGBColor(0x93,0xC5,0xFD), bold=True)
add_text_box(s, 3.8, 3.95, 2.5, 0.3, "Abhishek Sharma", font_size=16, color=WHITE, bold=True)
add_text_box(s, 3.8, 4.35, 2.5, 0.3, "UT-231-049-0004", font_size=14, color=WHITE)
add_text_box(s, 6.8, 3.6, 2.8, 0.3, "GUIDED BY", font_size=10, color=RGBColor(0x93,0xC5,0xFD), bold=True)
add_text_box(s, 6.8, 3.95, 2.8, 0.3, "Manas Chakraborty", font_size=16, color=WHITE, bold=True)
add_text_box(s, 6.8, 4.35, 2.8, 0.3, "Dept. of BCA", font_size=14, color=WHITE)
add_text_box(s, 1, 5.8, 11.3, 0.4, "ICON Commerce College  ·  2026", font_size=14, color=RGBColor(0xBA,0xE6,0xFD), alignment=PP_ALIGN.CENTER)

# ── SLIDE 2: Introduction ──
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
slide_header(s, "Introduction")
add_text_box(s, 0.8, 1.4, 11.7, 0.8, "ICC Companion is a web-based student portal built for ICON Commerce College. It digitizes daily academic operations — from attendance tracking to resource sharing — into a single unified platform for students, teachers, and administrators.", font_size=15, color=GRAY)
add_card(s, 0.8, 2.5, 5.6, 2.2, "Objective", ["Replace manual, paper-based processes", "with a real-time digital system that improves", "transparency and efficiency."], "🎯")
add_card(s, 6.7, 2.5, 5.6, 2.2, "Target Users", ["Students (BCA, BBA, BA, B.Com),", "Faculty members who manage classes,", "and College administrators."], "👥")

# ── SLIDE 3: Problem Statement ──
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, DARK)
slide_header(s, "Problem Statement", dark=True)
problems = ["→  Attendance tracked via paper registers — prone to errors and delays.",
            "→  Important notices shared through WhatsApp groups — easily missed.",
            "→  No centralized place for study materials, syllabus, or question papers.",
            "→  Students unaware of their attendance percentage until semester end."]
add_text_box(s, 0.8, 1.4, 5, 0.4, "Existing System Issues", font_size=18, bold=True, color=RGBColor(0x93,0xC5,0xFD))
y = 1.9
for p in problems:
    add_text_box(s, 0.8, y, 5.5, 0.35, p, font_size=13, color=RGBColor(0xE2,0xE8,0xF0))
    y += 0.42
prob_card = add_rect(s, 6.8, 1.4, 5.5, 3.0, RGBColor(0x33,0x41,0x55), RGBColor(0x47,0x55,0x69))
add_text_box(s, 7.1, 1.6, 5, 0.35, "Problem Definition", font_size=18, bold=True, color=WHITE)
add_text_box(s, 7.1, 2.1, 4.9, 2.0, "\"To design and develop a secure, user-friendly web portal that automates attendance tracking, centralizes academic resources, and provides real-time analytics for students and faculty at ICON Commerce College.\"", font_size=13, color=RGBColor(0xCB,0xD5,0xE1))

# ── SLIDE 4: Proposed System ──
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
slide_header(s, "Proposed System & Objectives")
add_text_box(s, 0.8, 1.4, 11, 0.4, "ICC Companion provides the following core features across Student, Admin, and Principal portals:", font_size=14, color=GRAY)
features = [
    ("📊", "Attendance Tracking", "Subject-wise real-time attendance\nwith 75% threshold alerts.", PRIMARY),
    ("📚", "Resource Library", "Upload & download syllabus, notes,\nbooks filtered by dept & semester.", WARNING),
    ("📅", "Exam Schedules", "View sessional and final exam\ntimetables with room details.", SECONDARY),
    ("📋", "Class Routines", "Department-wise weekly class\nroutine PDFs by administrators.", RGBColor(0x14,0xB8,0xA6)),
    ("🔍", "Lost & Found", "Report lost items or post found\nitems with claim status tracking.", DANGER),
    ("📢", "Announcements", "Priority-based notices targeted\nby department and semester.", PURPLE),
]
for i, (icon, title, desc, color) in enumerate(features):
    col = i % 3
    row = i // 3
    x = 0.8 + col * 4.1
    y_pos = 2.0 + row * 2.3
    add_card(s, x, y_pos, 3.8, 2.0, title, desc.split('\n'), icon, accent_color=color)

# ── SLIDE 5: System Modules ──
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
slide_header(s, "System Modules")
student_feats = ["Login with Roll Number & DOB", "View subject-wise attendance", "Download notes & resources", "View exam timetables", "Check class routines", "Browse announcements", "Report/view lost & found items"]
admin_feats = ["3 roles: Faculty · HOD · Principal", "Secure username/password login", "Mark daily attendance", "Manage students & resources", "HOD: Full department access", "Principal: System-wide access"]
super_feats = ["All HOD & Faculty privileges", "Add / Edit / Delete faculty accounts", "Manage departments", "Manage exams, routines, announcements", "Assign subjects to faculty", "Full CRUD over all data"]
add_card(s, 0.8, 1.4, 3.8, 4.5, "Student Portal", student_feats, "👨‍🎓")
add_card(s, 4.85, 1.4, 3.8, 4.5, "Admin Portal", admin_feats, "👨‍🏫")
add_card(s, 8.9, 1.4, 3.8, 4.5, "Principal (Super Admin)", super_feats, "👑")

# ── SLIDE 6: Technology Stack ──
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
slide_header(s, "Technology Stack")
tech = [
    ("Frontend", "HTML5, CSS3 (Custom Properties, Flexbox, Grid), JavaScript (ES6)"),
    ("Backend", "PHP (Server-side scripting via XAMPP Apache)"),
    ("Database", "MySQL (managed through phpMyAdmin via XAMPP)"),
    ("Local Server", "XAMPP (Apache + MySQL + PHP bundle)"),
    ("Design", "Mobile-first responsive layout, CSS variables for theming"),
]
# Table header
hdr = add_rect(s, 0.8, 1.5, 11.7, 0.5, PRIMARY)
add_text_box(s, 1.0, 1.52, 3, 0.4, "Layer", font_size=14, bold=True, color=WHITE)
add_text_box(s, 4.2, 1.52, 8, 0.4, "Technologies", font_size=14, bold=True, color=WHITE)
for i, (layer, techs) in enumerate(tech):
    y_pos = 2.05 + i * 0.55
    bg_color = LIGHT if i % 2 == 0 else WHITE
    row = add_rect(s, 0.8, y_pos, 11.7, 0.5, bg_color, BORDER)
    add_text_box(s, 1.0, y_pos + 0.02, 3, 0.4, layer, font_size=13, bold=True, color=DARK)
    add_text_box(s, 4.2, y_pos + 0.02, 8, 0.4, techs, font_size=13, color=DARK)

# ── SLIDE 7: System Architecture ──
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, DARK)
slide_header(s, "System Architecture", dark=True)
add_text_box(s, 0.8, 1.4, 11, 0.5, "The application follows a 3-Tier Architecture (Client → Server → Database).", font_size=15, color=RGBColor(0xCB,0xD5,0xE1))
# Arch boxes
arch_data = [
    ("🖥️ Presentation", "HTML · CSS · JS", PRIMARY, 1.5),
    ("⚙️ Application", "PHP · XAMPP Apache", WARNING, 5.5),
    ("🗄️ Data", "MySQL · phpMyAdmin", SECONDARY, 9.5),
]
for label, sub, color, x in arch_data:
    box = add_rect(s, x, 2.2, 2.8, 1.2, color)
    add_text_box(s, x + 0.1, 2.3, 2.6, 0.5, label, font_size=15, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x + 0.1, 2.8, 2.6, 0.4, sub, font_size=11, color=WHITE, alignment=PP_ALIGN.CENTER)
# Arrows
add_text_box(s, 4.3, 2.5, 1.2, 0.5, "⟷", font_size=28, color=GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(s, 8.3, 2.5, 1.2, 0.5, "⟷", font_size=28, color=GRAY, alignment=PP_ALIGN.CENTER)
# Bullets
arch_points = [
    "→  Client: Responsive web pages served to any browser (mobile or desktop).",
    "→  Server: PHP handles authentication, CRUD operations, and file uploads.",
    "→  Database: MySQL stores students, attendance, resources, announcements, and more.",
]
y = 3.8
for pt in arch_points:
    add_text_box(s, 0.8, y, 11, 0.35, pt, font_size=14, color=RGBColor(0xE2,0xE8,0xF0))
    y += 0.45

# ── SLIDE 8: Database Design ──
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
slide_header(s, "Database Design")
add_text_box(s, 0.8, 1.4, 11, 0.4, "The campus_portal database contains 11 normalized tables with foreign key constraints + 1 view.", font_size=14, color=GRAY)
tables = [
    ("departments", "dept_id · dept_code · dept_name", PRIMARY),
    ("students", "student_id · roll_number · name · dob", PRIMARY),
    ("admins", "admin_id · username · password · role", PRIMARY),
    ("subjects", "subject_id · code · name · dept · sem", PRIMARY),
    ("attendance", "id · roll_number(FK) · subject_id(FK)", SECONDARY),
    ("resources", "id · title · type · file_url · dept", SECONDARY),
    ("announcements", "id · title · priority · target_dept", SECONDARY),
    ("exams", "id · subject_id(FK) · date · room · type", SECONDARY),
    ("exam_schedules", "id · type (sessional/final) · file_url", WARNING),
    ("class_routines", "id · title · file_url · dept · semester", WARNING),
    ("lost_found", "id · title · category · type · status", WARNING),
    ("attendance_summary", "VIEW — pre-computed percentages", GRAY),
]
for i, (name, cols, color) in enumerate(tables):
    col = i % 4
    row = i // 4
    x = 0.8 + col * 3.1
    y_pos = 1.9 + row * 1.5
    card = add_rect(s, x, y_pos, 2.9, 1.2, WHITE, BORDER)
    add_text_box(s, x + 0.15, y_pos + 0.1, 2.6, 0.3, name, font_size=12, bold=True, color=color)
    add_text_box(s, x + 0.15, y_pos + 0.45, 2.6, 0.6, cols, font_size=9, color=GRAY)

# ── SLIDE 9: Security ──
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, DARK)
slide_header(s, "Security Features", dark=True)
sec_left = [
    "→  Role-Based Access: Students, Faculty, HOD, and Principal each have separate permissions.",
    "→  Session Management: PHP sessions to maintain login state and prevent unauthorized access.",
    "→  Input Validation: Server-side validation on all forms to prevent SQL injection and XSS.",
    "→  Secure File Uploads: File type and size restrictions on PDF and image uploads.",
]
sec_right = [
    "→  Password Protection: Admin credentials stored securely in the database.",
    "→  Student Auth: Login via Roll Number + Date of Birth — no complex passwords needed.",
    "→  Unique Constraints: Database-level enforcement to prevent duplicate attendance entries.",
]
y = 1.5
for pt in sec_left:
    add_text_box(s, 0.8, y, 6, 0.4, pt, font_size=13, color=RGBColor(0xE2,0xE8,0xF0))
    y += 0.5
y = 1.5
for pt in sec_right:
    add_text_box(s, 7.0, y, 5.5, 0.4, pt, font_size=13, color=RGBColor(0xE2,0xE8,0xF0))
    y += 0.5

# ── SLIDE 10: Testing ──
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
slide_header(s, "Testing")
tests = [
    ("Unit Testing", "Tested individual modules — login validation, attendance percentage calculation, PDF upload handling, and CRUD operations on all tables."),
    ("Integration Testing", "Verified end-to-end flows: admin marks attendance → database updates → student dashboard reflects the change instantly."),
    ("UI / Responsive Testing", "Tested across Chrome, Firefox, and mobile browsers. Ensured all pages render correctly on phones, tablets, and desktops."),
]
for i, (title, desc) in enumerate(tests):
    x = 0.8 + i * 4.1
    add_card(s, x, 1.5, 3.8, 2.8, title, [desc], accent_color=PRIMARY)

# ── SLIDE 11: Limitations & Future Scope ──
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
slide_header(s, "Limitations & Future Scope")
lim = ["Requires active internet/network connection.", "Currently runs on a local XAMPP server setup.", "No email or SMS notification system."]
fut = ["Deploy to a live web server for remote access.", "Online fee payment gateway integration.", "Email notifications for announcements.", "Automated end-of-semester report generation."]
# Limitations card
add_rect(s, 0.8, 1.5, 5.6, 3.5, WHITE, BORDER)
accent_l = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Emu(Inches(0.04)), Inches(3.5))
accent_l.fill.solid(); accent_l.fill.fore_color.rgb = DANGER; accent_l.line.fill.background()
add_text_box(s, 1.1, 1.7, 5, 0.35, "Current Limitations", font_size=18, bold=True, color=DANGER)
y = 2.2
for l in lim:
    add_text_box(s, 1.1, y, 5, 0.3, "•  " + l, font_size=13, color=GRAY)
    y += 0.4
# Future card
add_rect(s, 6.8, 1.5, 5.6, 3.5, WHITE, BORDER)
accent_r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.5), Emu(Inches(0.04)), Inches(3.5))
accent_r.fill.solid(); accent_r.fill.fore_color.rgb = SECONDARY; accent_r.line.fill.background()
add_text_box(s, 7.1, 1.7, 5, 0.35, "Future Enhancements", font_size=18, bold=True, color=SECONDARY)
y = 2.2
for f in fut:
    add_text_box(s, 7.1, y, 5, 0.3, "•  " + f, font_size=13, color=GRAY)
    y += 0.4

# ── SLIDE 12: Conclusion ──
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, DARK)
add_text_box(s, 1, 1.5, 11.3, 0.6, "Conclusion", font_size=32, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(s, 2, 2.3, 9.3, 1.2, "ICC Companion successfully digitizes the academic workflow at ICON Commerce College — replacing manual attendance registers, scattered WhatsApp notices, and disorganized resources with a single, transparent, real-time platform for students, faculty, and administrators.", font_size=16, color=RGBColor(0xCB,0xD5,0xE1), alignment=PP_ALIGN.CENTER)
add_text_box(s, 1, 3.8, 11.3, 0.8, "Thank You", font_size=48, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(s, 1, 4.7, 11.3, 0.5, "Open for Questions", font_size=18, color=RGBColor(0xBA,0xE6,0xFD), alignment=PP_ALIGN.CENTER)
add_text_box(s, 1, 5.5, 11.3, 0.4, "ICC Companion  ·  ICON Commerce College", font_size=13, color=RGBColor(0x94,0xA3,0xB8), alignment=PP_ALIGN.CENTER)

# Save
output = r"c:\xampp\htdocs\ICC_Companion\ICC_Companion_Presentation.pptx"
prs.save(output)
print(f"PPTX saved to: {output}")
